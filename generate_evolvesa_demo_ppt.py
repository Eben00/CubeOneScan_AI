from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(r"c:\AI\EvolveSA_Demo_Deck.pptx")
LOGO_PATH = Path(
    r"C:\Users\franc\.cursor\projects\c-AI\assets\c__Users_franc_AppData_Roaming_Cursor_User_workspaceStorage_68802244423911ff43916ecc5c9e8a12_images_EvolveSA_Photo-377c82a3-c5a3-4472-8fe3-1f0ce9b6400b.png"
)

GREEN = RGBColor(57, 255, 20)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
LIGHT_BG = RGBColor(245, 247, 250)


def add_header(slide, title):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    tf = header.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.8), Inches(0.08), height=Inches(0.72))


def add_bullets(slide, items, top=1.3):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(5.6))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24 if i == 0 else 20)
        p.font.bold = i == 0
        p.font.color.rgb = DARK if i == 0 else MUTED
        p.space_after = Pt(12)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(1.2), Inches(0.8), width=Inches(10.8))

    title = slide.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(11.3), Inches(1.8))
    tf = title.text_frame
    tf.text = "EvolveSA White-Label Demo"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GREEN

    sub = tf.add_paragraph()
    sub.text = "CubeOneScan platform customized for EvolveSA"
    sub.font.size = Pt(24)
    sub.font.color.rgb = WHITE

    footer = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11), Inches(0.6))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Prepared for stakeholder demo"
    fp.font.size = Pt(14)
    fp.font.color.rgb = RGBColor(156, 163, 175)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "1. Demo Goals")
    add_bullets(
        slide,
        [
            "What we are demonstrating today",
            "EvolveSA branding fully integrated in Android app flavor",
            "Launcher icon, login/home branding, and support metadata customized",
            "Working EvolveSA debug APK built and ready to install",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "2. White-Label Deliverables Completed")
    add_bullets(
        slide,
        [
            "Completed implementation",
            "Flavor-specific app name and strings (`EvolveSAScan`)",
            "EvolveSA logo integrated for in-app brand assets",
            "Flavor-specific launcher icon override via manifest merge",
            "EvolveSA-specific support and privacy details configured",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "3. Platform Flow")
    add_bullets(
        slide,
        [
            "User journey in the demo",
            "Login -> Home dashboard -> Scan workflow -> Lead creation",
            "Connector/API communication validated for tenant flow",
            "Operational screens accessible (settings, approvals, dashboard)",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "4. Build & Release Readiness")
    add_bullets(
        slide,
        [
            "Current build status",
            "EvolveSA flavor resources and manifest processing successful",
            "Full `assembleEvolvesaDebug` build successful",
            "Demo APK ready: `app-evolvesa-debug.apk`",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "5. Live Demo Script (3 Minutes)")
    add_bullets(
        slide,
        [
            "Suggested run-of-show",
            "Open app -> highlight EvolveSA icon and login branding",
            "Show settings + successful connection test",
            "Run one happy path scan/lead action and show response",
            "Close with deployment readiness and next phase",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "6. Next Steps After Demo")
    add_bullets(
        slide,
        [
            "Post-demo execution plan",
            "Confirm UAT feedback and required copy changes",
            "Finalize production signing and release pipeline",
            "Prepare Play Store/internal distribution package",
            "Schedule pilot rollout with support SLAs",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Thank You")
    add_bullets(
        slide,
        [
            "Questions & Discussion",
            "EvolveSA x CubeOne partnership demo",
        ],
        top=2.4,
    )

    prs.save(str(OUTPUT_PATH))
    print(f"Created: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
